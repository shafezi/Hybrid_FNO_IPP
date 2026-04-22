"""Generate talk.pptx for a 7-minute video presentation.

Reuses the visual style of Proposal_ECE568.pptx (UMich-flavored 16:9 layouts)
and populates 8 slides covering the full paper for an ICRA-style talk.
"""
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "Proposal_ECE568.pptx")
OUT = os.path.join(ROOT, "talk.pptx")
FIG = os.path.join(ROOT, "results", "dynamic_ipp", "final")

# UMich palette
UMBLUE  = RGBColor(0x00, 0x27, 0x4C)
UMMAIZE = RGBColor(0xFF, 0xCB, 0x05)
TEXT    = RGBColor(0x22, 0x22, 0x22)
GREY    = RGBColor(0x66, 0x66, 0x66)


# ---------------------------------------------------------------------------- #
def clear_slides(prs):
    """Remove all existing slides from a presentation (XML-level)."""
    sldIdLst = prs.slides._sldIdLst
    rIds = [s.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
            for s in sldIdLst]
    for slide in list(sldIdLst):
        sldIdLst.remove(slide)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def get_layout(prs, name):
    """Find a slide layout by name across all masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise ValueError(f"No layout named {name!r}")


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = 'Calibri'
    return tb


def add_bullets(slide, x, y, w, h, bullets, size=16, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        # Add bullet character manually since we're not using a placeholder
        r = p.add_run()
        r.text = u"•  " + b
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = 'Calibri'
        p.space_after = Pt(6)
    return tb


def add_title_bar(slide, prs, title, subtitle=None):
    """Solid umblue bar across the top with title (white) and optional subtitle."""
    sw = prs.slide_width
    bar_h = Inches(1.1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = UMBLUE
    bar.line.fill.background()
    # Maize accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, bar_h, sw, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = UMMAIZE
    stripe.line.fill.background()
    # Title text
    add_textbox(slide, Inches(0.5), Inches(0.18), sw - Inches(1.0), Inches(0.55),
                title, size=28, bold=True, color=RGBColor(255, 255, 255))
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.70), sw - Inches(1.0), Inches(0.35),
                    subtitle, size=14, italic=True, color=RGBColor(220, 220, 220))


def add_footer(slide, prs, page_no, total):
    sw, sh = prs.slide_width, prs.slide_height
    add_textbox(slide, Inches(0.4), sh - Inches(0.4), Inches(6.0), Inches(0.3),
                "A Modular Framework for Multi-Robot Adaptive Ocean Monitoring  |  S. Hafezi  |  UMich",
                size=10, color=GREY)
    add_textbox(slide, sw - Inches(1.0), sh - Inches(0.4), Inches(0.7), Inches(0.3),
                f"{page_no} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------- #
def set_notes(slide, text):
    """Attach speaker notes to a slide."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ---------------------------------------------------------------------------- #
NOTES = {
1: """Hi, I'm Shayesteh Hafezi from the Department of Naval Architecture and Marine Engineering at Michigan. I'd like to talk about a framework I built for multi-robot adaptive ocean monitoring.""",

2: """So here's the problem.

If you want to monitor the ocean — track eddies, watch fronts evolve — you've got a tough scale mismatch. Underwater gliders sample sub-kilometer features really well, but only along the line they travel. And you typically only have a handful of them, trying to cover ocean basins thousands of kilometers wide.

So you need to combine sparse mobile observations with some kind of forecast.

Now, two research communities have worked on parts of this, but separately. The classical informative-path-planning folks fit a Gaussian process over the field and route robots to where they'll learn the most. The forecast folks build learned models — FNO, OceanNet — that produce great basin-scale predictions, but offline, with no in-situ data being fed in. Nobody has closed the loop between them. That's the gap I'm trying to fill.""",

3: """Here's the framework. It's a closed loop with five steps.

Step one: a dynamical prior makes a forecast. That could be a learned neural operator like FNO, or it could be something trivial like persistence — just hold the field constant. Whatever you've got.

Step two: a Gaussian process fits the residual. So instead of fitting the full field, the GP only models where the prior is wrong, based on the observations the robots have collected so far.

Step three: each robot picks its next sample location — but only within its own Voronoi cell. That's how I keep two robots from going after the same target.

Step four: robots actually go there, take a measurement, the GP updates.

Step five: the corrected estimate — prior plus GP correction — feeds back into step one, and we loop.

The important thing is that f, the prior, is a black box. The same code runs whether I plug in FNO or persistence — it's genuinely modular.""",

4: """OK so here's the most interesting thing I found, and honestly it surprised me.

While I was benchmarking the framework, I noticed something weird about the metrics.

Look on the left — these are estimated ocean fields at day 40 with 20 robots. The ground truth has all this fine structure: eddies, fronts. The methods that use a prior — FNO+GP, Persist+GP — they capture that structure. But the no-prior method, just a GP fit to the observations, looks like a smooth blob.

Now here's the catch. That smooth blob has the lowest RMSE of any method. The reason is, RMSE rewards anything close to the field's average value — and a smooth blob basically is the average.

So I computed power spectra — these are 1D spectra over ocean transects on the right — and you can see the GP-only method collapses by four to five orders of magnitude at high wavenumbers. It just has no fine-scale content. The structural information is completely gone. And RMSE doesn't see any of this.

So my methodology argument is: the IPP community should adopt spectral evaluation alongside RMSE. We're systematically misranking methods otherwise.""",

5: """Here's the same point in numbers. This is at 20 robots — RMSE, ACC, HF, and FSS for all five methods.

If you only look at RMSE, the leftmost panel — the GP-only method wins. RMSE of 0.78. By any standard publication, you'd say it's the best method.

But look at HF, the high-frequency energy ratio, third panel. The ideal value is 1 — meaning your prediction has the same fine-scale energy as the truth. GP-only gets 0.10 — basically nothing. FNO+GP gets 0.90, very close to ideal.

So the rankings flip completely depending on which metric you read. Read RMSE alone, GP-only wins. Read HF alone, FNO+GP wins. Only the joint reading is faithful to what's actually happening to the field.""",

6: """To make sure these findings aren't just a hyperparameter artifact, I ran the full sweep — 4 kernels, 3 acquisition functions, 4 robot counts, 5 random seeds. About 1,200 episodes in total.

On the left is the Pareto plot. Every dot is one configuration. You can see the methods occupy completely distinct regions in metric space — the colors don't overlap. So whatever kernel or acquisition function you pick, the ranking holds.

On the right is how things scale with robot count. The benefit of the GP correction loop grows with how many robots you have. Below about 10 robots, all methods struggle. Above 10, FNO+GP clearly outperforms uncorrected FNO across every metric.""",

7: """So to wrap up — three contributions.

First, the modular closed-loop framework. Any prior, plus GP residual correction, plus multi-robot Voronoi acquisition. The same code runs with FNO, persistence, or no-prior. That's the modularity claim.

Second — and I think this is actually the more important contribution — a spectral evaluation methodology for IPP. The 1D Hanning-windowed PSD on ocean transects, plus structural metrics like ACC and FSS. These tools are standard in weather forecasting, but the IPP community doesn't use them. We should.

And third, the empirical analysis. The GP loop helps any prior. Learned priors give consistently better structural metrics than naive ones. And methods without any dynamical prior just can't recover fine-scale ocean structure, regardless of how you tune them.""",

8: """There's a lot of future work I'd like to do. Multi-year and multi-variable generalization, hardware-in-the-loop validation with real glider fleets, and multi-objective acquisition for tracking specific features like eddies and biological hotspots.

Thanks to the ROMS team for the regional ocean reanalyses, and the OceanNet authors for releasing their model weights.

And thank you for watching.""",
}


def build():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # Use a generic blank layout for all slides; we draw our own headers.
    blank = get_layout(prs, "Content Slide 2")

    sw, sh = prs.slide_width, prs.slide_height
    TOTAL = 8

    # -------- SLIDE 1 : Title -------------------------------------------------
    s = prs.slides.add_slide(blank)
    # Full-bleed umblue background band on top half
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    bg.line.fill.background()
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), sw, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = UMBLUE
    band.line.fill.background()
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.4), sw, Inches(0.10))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = UMMAIZE
    stripe.line.fill.background()

    add_textbox(s, Inches(0.6), Inches(2.5), sw - Inches(1.2), Inches(1.5),
                "A Modular Framework for\nMulti-Robot Adaptive Ocean Monitoring",
                size=40, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(4.3), sw - Inches(1.2), Inches(0.6),
                "with Gaussian-Process Residual Correction",
                size=22, italic=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.7), sw - Inches(1.2), Inches(0.4),
                "Shayesteh Hafezi",
                size=20, bold=True, color=UMBLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.1), sw - Inches(1.2), Inches(0.4),
                "Department of Naval Architecture and Marine Engineering, University of Michigan",
                size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.5), sw - Inches(1.2), Inches(0.4),
                "shafezi@umich.edu  |  April 2026",
                size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    set_notes(s, NOTES[1])

    # -------- SLIDE 2 : Problem & Gap -----------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "The Problem", "Multi-robot ocean monitoring needs both forecasts and adaptive observations")
    add_bullets(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.0), [
        "Underwater gliders sample sub-kilometer ocean features in situ",
        "But a handful of platforms must cover ocean basins spanning thousands of km",
        "",
        "Two research communities have grown around this problem — separately:",
        "    Classical IPP: a GP over the field, no physics",
        "    Learned forecasts (FNO, OceanNet): basin-scale, no in-situ data",
        "",
        "Gap: no framework integrates a learned prior with multi-robot adaptive sampling",
    ], size=18)
    # Right-side visual: ocean schematic placeholder using the system diagram thumbnail
    img = os.path.join(FIG, "system_diagram.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(7.3), Inches(2.2), width=Inches(5.8))
    add_footer(s, prs, 2, TOTAL)
    set_notes(s, NOTES[2])

    # -------- SLIDE 3 : Closed-Loop Framework ---------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "A Modular Closed-Loop Framework",
                  "Any dynamical prior + GP residual correction + multi-robot Voronoi acquisition")
    img = os.path.join(FIG, "system_diagram.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.4), Inches(1.5), width=Inches(8.4))
    add_bullets(s, Inches(9.0), Inches(1.7), Inches(4.0), Inches(5.0), [
        "Prior  ŷ_prior = f(ŷ)",
        "GP fits residual  e = y_true − ŷ_prior",
        "Voronoi acquisition picks next sites",
        "Robots execute (0.5 m/s glider speed)",
        "Corrected: ŷ = ŷ_prior + μ_GP",
        "",
        "Same code runs FNO, persistence,",
        "or no-prior — f is a black box.",
    ], size=14)
    add_footer(s, prs, 3, TOTAL)
    set_notes(s, NOTES[3])

    # -------- SLIDE 4 : Methodology Insight (the punchline) -------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "RMSE Is Systematically Misleading",
                  "A no-prior GP wins RMSE — by predicting the mean")
    img1 = os.path.join(FIG, "figures",
                        "seed42_t030_20bots_matern15_uncertainty_only_step10.png")
    img2 = os.path.join(FIG, "figures_paper", "power_spectra.png")
    if os.path.exists(img1):
        s.shapes.add_picture(img1, Inches(0.3), Inches(1.5), width=Inches(7.4))
    if os.path.exists(img2):
        s.shapes.add_picture(img2, Inches(7.9), Inches(1.5), width=Inches(5.2))
    add_textbox(s, Inches(0.3), Inches(5.9), Inches(7.4), Inches(0.4),
                "Qualitative: GP-only is a smooth blob",
                size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.9), Inches(5.9), Inches(5.2), Inches(0.4),
                "PSD: GP-only collapses 4–5 orders of magnitude at high k",
                size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.8), [
        "We adopt 1D Hanning-windowed PSD (standard SSH protocol) and structural metrics (ACC, FSS) as the proper evaluation tools for IPP."
    ], size=13)
    add_footer(s, prs, 4, TOTAL)
    set_notes(s, NOTES[4])

    # -------- SLIDE 5 : Quantitative Results ----------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Per-Method Comparison",
                  "Read RMSE alone: GP-only wins. Read HF alone: FNO+GP wins.")
    img = os.path.join(FIG, "figures_paper", "bar_chart_20bots.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.4), Inches(1.5), width=Inches(9.0))
    add_textbox(s, Inches(9.6), Inches(1.6), Inches(3.5), Inches(0.5),
                "n = 20 robots, day 40", size=14, bold=True, color=UMBLUE)
    add_bullets(s, Inches(9.6), Inches(2.1), Inches(3.5), Inches(4.0), [
        "GP-only:    RMSE 0.78,  HF 0.10",
        "FNO+GP:   RMSE 0.94,  HF 0.90",
        "Persist+GP: RMSE 1.15, HF 0.85",
        "",
        "GP-only has the lowest RMSE",
        "but ≈0 high-frequency content.",
        "",
        "Only the joint reading is faithful.",
    ], size=13)
    add_footer(s, prs, 5, TOTAL)
    set_notes(s, NOTES[5])

    # -------- SLIDE 6 : Robustness across hyperparameters / scaling ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Robustness and Scaling",
                  "4 kernels × 3 acquisitions × 4 robot counts × 5 seeds")
    pareto = os.path.join(FIG, "figures_paper", "pareto_grid_per_seed_plain.png")
    scaling = os.path.join(FIG, "figures_paper", "scaling_curve.png")
    if os.path.exists(pareto):
        s.shapes.add_picture(pareto, Inches(0.3), Inches(1.4), height=Inches(5.5))
    if os.path.exists(scaling):
        s.shapes.add_picture(scaling, Inches(6.5), Inches(1.5), width=Inches(6.7))
    add_bullets(s, Inches(6.5), Inches(5.3), Inches(6.7), Inches(1.5), [
        "Methods occupy distinct metric regions; rankings robust to kernel and acquisition",
        "Framework benefit grows with observation density; minimum useful regime ≈ 10 robots",
    ], size=13)
    add_footer(s, prs, 6, TOTAL)
    set_notes(s, NOTES[6])

    # -------- SLIDE 7 : Take-aways --------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Take-Aways",
                  "Three contributions of this work")
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.5), Inches(0.5),
                "(C1)  A modular closed-loop framework",
                size=22, bold=True, color=UMBLUE)
    add_textbox(s, Inches(0.9), Inches(2.05), Inches(12.0), Inches(0.6),
                "Any prior + GP residual correction + multi-robot Voronoi acquisition. The same code runs with FNO, persistence, or no-prior.",
                size=15, color=TEXT)

    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12.5), Inches(0.5),
                "(C2)  A spectral evaluation methodology for IPP",
                size=22, bold=True, color=UMBLUE)
    add_textbox(s, Inches(0.9), Inches(3.55), Inches(12.0), Inches(0.6),
                "1D Hanning-windowed PSD on ocean-only transects + structural metrics. Exposes RMSE failure modes in IPP.",
                size=15, color=TEXT)

    add_textbox(s, Inches(0.5), Inches(4.5), Inches(12.5), Inches(0.5),
                "(C3)  Empirical analysis across the full hyperparameter sweep",
                size=22, bold=True, color=UMBLUE)
    add_textbox(s, Inches(0.9), Inches(5.05), Inches(12.0), Inches(1.0),
                "GP loop helps any prior; learned priors yield consistently better structural metrics; no-prior methods cannot recover fine-scale ocean structure regardless of tuning.",
                size=15, color=TEXT)

    add_footer(s, prs, 7, TOTAL)
    set_notes(s, NOTES[7])

    # -------- SLIDE 8 : Future Work / Thank You -------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Future Work and Acknowledgments",
                  "")
    add_textbox(s, Inches(0.5), Inches(1.6), Inches(12.5), Inches(0.5),
                "Future Work", size=22, bold=True, color=UMBLUE)
    add_bullets(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(2.0), [
        "Multi-year and multi-variable generalization (beyond NW Pacific SSH 2020)",
        "Hardware-in-the-loop validation with real glider fleets",
        "Multi-objective acquisition for feature tracking (eddies, fronts, blooms)",
        "Stronger priors in the sparse regime (n ≤ 5 robots)",
    ], size=15)

    add_textbox(s, Inches(0.5), Inches(4.6), Inches(12.5), Inches(0.5),
                "Acknowledgments", size=22, bold=True, color=UMBLUE)
    add_textbox(s, Inches(0.7), Inches(5.1), Inches(12.0), Inches(0.6),
                "ROMS team for regional ocean reanalyses; OceanNet authors for releasing model weights.",
                size=14, color=TEXT)

    # Big "Thank You" + contact
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.2), sw, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = UMBLUE; band.line.fill.background()
    add_textbox(s, Inches(0.5), Inches(6.35), sw - Inches(1.0), Inches(0.8),
                "Thank you  •  shafezi@umich.edu",
                size=22, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    add_footer(s, prs, 8, TOTAL)
    set_notes(s, NOTES[8])

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Open in PowerPoint / Keynote / LibreOffice Impress.")

    # Also write a standalone speaker-script markdown file for easy reference
    script_path = os.path.join(ROOT, "talk_script.md")
    titles = {
        1: "Title",
        2: "The Problem",
        3: "A Modular Closed-Loop Framework",
        4: "RMSE Is Systematically Misleading",
        5: "Per-Method Comparison",
        6: "Robustness and Scaling",
        7: "Take-Aways",
        8: "Future Work and Acknowledgments",
    }
    times = {1: "10s", 2: "55s", 3: "65s", 4: "85s", 5: "55s", 6: "55s", 7: "60s", 8: "30s"}
    with open(script_path, "w") as f:
        f.write("# Talk script — 7-minute video\n\n")
        f.write("Speaker notes for each slide. Also embedded inside `talk.pptx` "
                "(visible in PowerPoint's presenter view).\n\n")
        f.write("**Total estimated runtime: ~6:50**\n\n---\n\n")
        for i in range(1, 9):
            f.write(f"## Slide {i} — {titles[i]}  *(~{times[i]})*\n\n")
            f.write(NOTES[i].strip() + "\n\n---\n\n")
    print(f"Saved script: {script_path}")


if __name__ == "__main__":
    build()
